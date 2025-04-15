import os
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Utility functions used in the tool
def get_chart_type(chart_type_str):
    chart_types = {
        'bar-horizontal': XL_CHART_TYPE.BAR_CLUSTERED,
        'bar-vertical': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
    }
    return chart_types.get(chart_type_str.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)


def find_layout_by_name(prs, layout_name):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    return None

def run(message, params):
    """
    Generates a PowerPoint presentation based on the provided parameters.

    Args:
        message (str): Message or additional query input (not used in this tool).
        params (dict): Parameters including:
            - csv_path (str): Path to the input Excel/CSV file.
            - template_path (str): Path to the PowerPoint template file.

    Returns:
        dict: Status and message of the operation.
    """
    try:
        # Extract parameters
        csv_path = params.get("csv_path")
        template_path = params.get("template_path")

        # Validate input parameters
        if not csv_path or not os.path.exists(csv_path):
            return {"status": "error", "message": f"CSV file not found at '{csv_path}'."}
        if not template_path or not os.path.exists(template_path):
            return {"status": "error", "message": f"Template file not found at '{template_path}'."}

        # Load the dataset from CSV/Excel file
        df = pd.read_excel(csv_path)

        # Load the presentation template
        prs = Presentation(template_path)

        # Process each unique question_id
        for question_id in df['question_id'].unique():
            question_data = df[df['question_id'] == question_id]
            question_text = question_data['question_text'].iloc[0]
            text_summary = question_data['text_summary'].iloc[0]
            chart_type_str = question_data['chart_type'].iloc[0]
            chart_type = get_chart_type(chart_type_str)
            chart_layout_name = question_data['chart_layout'].iloc[0]

            # Find the slide layout by name and add a new slide
            slide_layout = find_layout_by_name(prs, chart_layout_name)
            slide = prs.slides.add_slide(slide_layout)

            # Set the title for the slide
            slide.shapes.title.text = question_text

            # Add chart
            chart_data = CategoryChartData()
            chart_data.categories = question_data['response'].tolist()
            chart_data.add_series('Series 1', (question_data['value'] * 100).tolist())  # Convert to percentage

            # Insert text_summary into the specified text placeholder
            text_placeholder = slide.placeholders[11]
            text_placeholder.text = str(text_summary)

            # Choose the chart placeholder by index and insert the chart
            chart_placeholder = slide.placeholders[10]
            chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
            chart = chart_frame.chart

            # Formatting based on the chart type
            if chart_type in (XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.LINE_MARKERS):
                for series in chart.series:
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(0x14, 0x60, 0x82)  # Default blue color

                    # Add data labels formatted as percentages with no decimal points
                    series.has_data_labels = True
                    for point in series.points:
                        point.data_label.number_format = '0%'
                        point.data_label.font.size = Pt(10)
                        point.data_label.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black font color

                # Remove the chart title and legend
                chart.has_title = False
                chart.has_legend = False

        # Save the presentation
        output_file = "output_presentation.pptx"
        prs.save(output_file)

        return {"status": "success", "message": f"Presentation created successfully! Saved as '{output_file}'."}

    except Exception as e:
        return {"status": "error", "message": str(e)}
